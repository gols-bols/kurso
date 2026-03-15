from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path("/Users/gnome/Downloads/kurso-main 2")
LIB_DIR = Path("/tmp/pptxlib")
if str(LIB_DIR) not in sys.path:
    sys.path.insert(0, str(LIB_DIR))

from pptx import Presentation  # type: ignore
from pptx.dml.color import RGBColor  # type: ignore
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR  # type: ignore
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR  # type: ignore
from pptx.util import Inches, Pt  # type: ignore


OUTPUT = ROOT / "docs" / "Презентация_к_курсовой.pptx"

BG = RGBColor(247, 240, 231)
PANEL = RGBColor(255, 251, 245)
ACCENT = RGBColor(142, 74, 42)
ACCENT_SOFT = RGBColor(224, 196, 164)
TEXT = RGBColor(35, 27, 20)
MUTED = RGBColor(108, 92, 73)
LINE = RGBColor(220, 203, 182)
GREEN = RGBColor(115, 144, 123)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_header(slide, title: str, tag: str, index: int, total: int) -> None:
    set_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.55)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    tag_box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(10.55), Inches(0.82), Inches(2.2), Inches(0.42)
    )
    tag_box.fill.solid()
    tag_box.fill.fore_color.rgb = ACCENT_SOFT
    tag_box.line.fill.background()
    tf = tag_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = tag
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.95), Inches(10.8), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT

    footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.8), Inches(0.25))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Слайд {index} / {total}"
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.color.rgb = MUTED


def add_card(slide, left: float, top: float, width: float, height: float, title: str | None = None):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.2)

    if title:
        box = slide.shapes.add_textbox(Inches(left + 0.22), Inches(top + 0.14), Inches(width - 0.44), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.name = "Arial"
        run.font.size = Pt(18)
        run.font.bold = True
        run.font.color.rgb = ACCENT

    return shape


def add_bullets(slide, left: float, top: float, width: float, height: float, bullets: list[str], font_size: int = 20) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.clear()

    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.line_spacing = 1.15
        p.space_after = Pt(10)
        p.bullet = True


def add_paragraph(slide, left: float, top: float, width: float, height: float, text: str, size: int = 20, color=TEXT, bold=False, align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = ACCENT
    connector.line.width = Pt(2.2)
    connector.line.end_arrowhead = True


def add_box(slide, left: float, top: float, width: float, height: float, text: str, fill_rgb=ACCENT_SOFT, text_rgb=TEXT, font_size=18) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_rgb


def slide_title(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.72)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    add_paragraph(
        slide,
        0.8,
        1.2,
        11.4,
        1.5,
        "Разработка информационной системы управления заявками",
        size=30,
        bold=True,
    )
    add_paragraph(
        slide,
        0.82,
        2.35,
        10.8,
        1.2,
        "Курсовой проект по разработке web-приложения на Laravel с использованием MySQL и Git.",
        size=21,
        color=MUTED,
    )

    add_card(slide, 0.82, 4.2, 5.0, 1.45, "Основные технологии")
    add_bullets(slide, 1.05, 4.65, 4.55, 0.9, ["Laravel 11", "PHP 8.2", "MySQL", "Git"], font_size=19)

    add_card(slide, 6.2, 4.2, 6.1, 1.45, "Назначение проекта")
    add_paragraph(
        slide,
        6.45,
        4.7,
        5.5,
        0.7,
        "Централизованный учет, просмотр и сопровождение внутренних заявок пользователей.",
        size=19,
    )

    add_paragraph(slide, 0.82, 6.95, 12.0, 0.25, "Слайд 1 / 11", size=12, color=MUTED, align=PP_ALIGN.RIGHT)


def slide_agenda(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Актуальность, цель и задачи проекта", "Основа работы", 2, total)
    add_card(slide, 0.82, 1.7, 5.75, 4.75, "Актуальность")
    add_bullets(
        slide,
        1.06,
        2.15,
        5.2,
        3.9,
        [
            "Устные и разрозненные обращения легко теряются.",
            "Пользователь не видит текущий статус своей заявки.",
            "Руководителю сложно контролировать ход обработки обращений.",
            "Нужна единая цифровая система регистрации и сопровождения.",
        ],
        font_size=18,
    )
    add_card(slide, 6.8, 1.7, 5.7, 4.75, "Цель и задачи")
    add_bullets(
        slide,
        7.05,
        2.15,
        5.1,
        3.9,
        [
            "Разработать рабочий сайт для управления заявками.",
            "Спроектировать базу данных MySQL.",
            "Реализовать роли, вход, просмотр, создание и редактирование.",
            "Подготовить документацию и материалы для защиты.",
        ],
        font_size=18,
    )


def slide_stack(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Технологический стек", "Технологии", 3, total)
    add_card(slide, 0.82, 1.75, 4.0, 4.6, "Backend")
    add_bullets(slide, 1.05, 2.2, 3.5, 3.7, ["PHP 8.2", "Laravel 11", "Eloquent ORM", "Blade"], font_size=20)

    add_card(slide, 4.95, 1.75, 3.2, 4.6, "База данных")
    add_bullets(slide, 5.18, 2.2, 2.7, 3.7, ["MySQL", "Миграции", "Сидеры", "Сессии"], font_size=20)

    add_card(slide, 8.3, 1.75, 4.1, 4.6, "Разработка")
    add_bullets(slide, 8.55, 2.2, 3.5, 3.7, ["Git", "GitHub", "Composer", "Локальный запуск"], font_size=20)

    add_paragraph(
        slide,
        0.95,
        6.45,
        11.5,
        0.35,
        "React в проекте не использовался, так как система реализована как server-rendered web-приложение на Laravel и Blade.",
        size=16,
        color=MUTED,
    )


def slide_architecture(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Общая архитектура приложения", "Диаграмма 1", 4, total)

    add_box(slide, 0.8, 3.0, 2.0, 0.9, "Пользователь", fill_rgb=GREEN, text_rgb=RGBColor(255, 255, 255))
    add_box(slide, 3.2, 3.0, 2.0, 0.9, "Браузер", fill_rgb=ACCENT_SOFT)
    add_box(slide, 5.6, 2.55, 2.2, 1.8, "Laravel\nприложение", fill_rgb=ACCENT, text_rgb=RGBColor(255, 255, 255))
    add_box(slide, 8.3, 1.5, 2.0, 0.7, "Маршруты", fill_rgb=PANEL)
    add_box(slide, 8.3, 2.45, 2.0, 0.7, "Контроллеры", fill_rgb=PANEL)
    add_box(slide, 8.3, 3.4, 2.0, 0.7, "Blade views", fill_rgb=PANEL)
    add_box(slide, 8.3, 4.35, 2.0, 0.7, "Модели", fill_rgb=PANEL)
    add_box(slide, 10.9, 3.0, 1.7, 0.9, "MySQL", fill_rgb=GREEN, text_rgb=RGBColor(255, 255, 255))

    add_arrow(slide, 2.8, 3.45, 3.2, 3.45)
    add_arrow(slide, 5.2, 3.45, 5.6, 3.45)
    add_arrow(slide, 7.8, 3.45, 8.3, 3.45)
    add_arrow(slide, 10.3, 3.45, 10.9, 3.45)

    add_paragraph(
        slide,
        0.95,
        5.95,
        11.2,
        0.6,
        "Пользователь взаимодействует с системой через браузер. Запросы обрабатываются Laravel-приложением, а постоянное хранение данных выполняется в MySQL.",
        size=18,
        color=MUTED,
    )


def slide_mvc(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "MVC-логика проекта", "Диаграмма 2", 5, total)

    add_box(slide, 0.9, 2.8, 1.7, 0.85, "Маршрут")
    add_box(slide, 3.0, 2.8, 1.9, 0.85, "Контроллер")
    add_box(slide, 5.4, 2.8, 1.8, 0.85, "Модель")
    add_box(slide, 7.6, 2.8, 1.7, 0.85, "MySQL", fill_rgb=GREEN, text_rgb=RGBColor(255, 255, 255))
    add_box(slide, 9.7, 2.8, 2.0, 0.85, "Blade view")

    add_arrow(slide, 2.6, 3.23, 3.0, 3.23)
    add_arrow(slide, 4.9, 3.23, 5.4, 3.23)
    add_arrow(slide, 7.2, 3.23, 7.6, 3.23)
    add_arrow(slide, 9.3, 3.23, 9.7, 3.23)

    add_card(slide, 1.0, 4.6, 10.8, 1.45, "Смысл схемы")
    add_paragraph(
        slide,
        1.25,
        5.02,
        10.2,
        0.7,
        "Контроллер получает HTTP-запрос, валидирует данные, обращается к модели и передает результат в Blade-представление для отображения интерфейса.",
        size=18,
    )


def slide_db(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Структура базы данных MySQL", "Диаграмма 3", 6, total)

    add_box(slide, 1.0, 2.05, 3.3, 2.7, "users\n\nid\nname\nemail\npassword\nrole", fill_rgb=PANEL)
    add_box(slide, 5.0, 1.7, 3.6, 3.4, "tickets\n\nid\ntitle\ndescription\npriority\nstatus\ncreated_by", fill_rgb=ACCENT_SOFT)
    add_box(slide, 9.2, 2.2, 2.5, 2.3, "sessions\n\nid\nuser_id\npayload\nlast_activity", fill_rgb=PANEL)

    add_arrow(slide, 4.3, 3.4, 5.0, 3.0)
    add_arrow(slide, 8.6, 3.2, 9.2, 3.2)

    add_paragraph(
        slide,
        0.95,
        5.75,
        11.3,
        0.7,
        "Ключевая связь проекта: tickets.created_by → users.id. Таблица sessions обеспечивает стабильную работу авторизации через сессии Laravel.",
        size=18,
        color=MUTED,
    )


def slide_functionality(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Реализованный функционал и роли", "Практика", 7, total)
    add_card(slide, 0.82, 1.75, 5.9, 4.85, "Функциональные возможности")
    add_bullets(
        slide,
        1.08,
        2.2,
        5.3,
        4.0,
        [
            "Вход и выход из системы.",
            "Просмотр списка заявок.",
            "Создание заявки только ролью user.",
            "Редактирование существующих заявок администратором и менеджером.",
            "Тестовые данные через сидеры.",
        ],
        font_size=18,
    )

    add_card(slide, 6.95, 1.75, 5.35, 4.85, "Ролевая модель")
    add_bullets(
        slide,
        7.2,
        2.2,
        4.7,
        4.0,
        [
            "user — создает свои заявки и редактирует свои.",
            "admin — просматривает и редактирует существующие заявки.",
            "manager — просматривает и редактирует существующие заявки.",
        ],
        font_size=18,
    )


def slide_git(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Git и GitHub", "Контроль версий", 8, total)
    add_card(slide, 0.82, 1.8, 5.8, 4.8, "Зачем нужен Git")
    add_bullets(
        slide,
        1.06,
        2.22,
        5.2,
        4.0,
        [
            "Фиксация этапов разработки проекта.",
            "История изменений по файлам и функциональности.",
            "Подтверждение самостоятельной работы над курсовой.",
            "Подготовка проекта к размещению на GitHub.",
        ],
        font_size=18,
    )
    add_card(slide, 6.92, 1.8, 5.38, 4.8, "Репозиторий проекта")
    add_paragraph(
        slide,
        7.18,
        2.3,
        4.8,
        0.65,
        "GitHub: https://github.com/gols-bols/kurso",
        size=18,
        bold=True,
    )
    add_bullets(
        slide,
        7.18,
        3.0,
        4.8,
        2.8,
        [
            "Код сайта и документация.",
            "Миграции и сидеры базы данных.",
            "Отдельная ветка course-defense-final.",
            "Готовность к дальнейшим коммитам и pull request.",
        ],
        font_size=18,
    )


def slide_docs(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "Документация и материалы защиты", "Сдача проекта", 9, total)
    add_card(slide, 0.82, 1.8, 5.8, 4.8, "Подготовленные материалы")
    add_bullets(
        slide,
        1.06,
        2.22,
        5.2,
        4.0,
        [
            "Пояснительная записка, оформленная по дипломной методичке.",
            "Презентация для выступления на 5–10 минут.",
            "Рабочий сайт с ролями, БД и сидерами.",
            "Готовая структура для локального запуска и демонстрации.",
        ],
        font_size=18,
    )
    add_card(slide, 6.92, 1.8, 5.38, 4.8, "Что показать на защите")
    add_bullets(
        slide,
        7.18,
        2.22,
        4.8,
        4.0,
        [
            "Страницу входа.",
            "Список заявок.",
            "Создание заявки пользователем.",
            "Редактирование заявки администратором.",
            "Схему БД и архитектурные диаграммы.",
        ],
        font_size=18,
    )


def slide_plan(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, "План выступления на защите", "5–10 минут", 10, total)
    add_card(slide, 1.0, 1.85, 11.1, 4.9, "Рекомендуемая структура доклада")
    add_bullets(
        slide,
        1.3,
        2.35,
        10.3,
        4.0,
        [
            "1 минута — актуальность темы и цель проекта.",
            "1–2 минуты — выбранные технологии: Laravel, MySQL, Git.",
            "2 минуты — архитектура, база данных и роли.",
            "2 минуты — живая демонстрация сайта.",
            "1 минута — выводы и направления развития.",
        ],
        font_size=20,
    )


def slide_final(prs: Presentation, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.72)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    add_paragraph(slide, 1.0, 1.8, 10.8, 0.9, "Спасибо за внимание", size=32, bold=True)
    add_paragraph(
        slide,
        1.0,
        2.9,
        11.0,
        1.0,
        "Готов ответить на вопросы по архитектуре, базе данных, ролям пользователей, Git и дальнейшему развитию проекта.",
        size=22,
        color=MUTED,
    )

    add_card(slide, 2.0, 4.4, 9.1, 1.25, "Ключевой результат")
    add_paragraph(
        slide,
        2.35,
        4.85,
        8.4,
        0.45,
        "Создан рабочий web-сайт управления заявками на Laravel с БД MySQL и подготовлен полный комплект материалов к защите.",
        size=19,
        align=PP_ALIGN.CENTER,
    )

    add_paragraph(slide, 0.82, 6.95, 12.0, 0.25, "Слайд 11 / 11", size=12, color=MUTED, align=PP_ALIGN.RIGHT)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = 11
    slide_title(prs, total)
    slide_agenda(prs, total)
    slide_stack(prs, total)
    slide_architecture(prs, total)
    slide_mvc(prs, total)
    slide_db(prs, total)
    slide_functionality(prs, total)
    slide_git(prs, total)
    slide_docs(prs, total)
    slide_plan(prs, total)
    slide_final(prs, total)

    prs.save(OUTPUT)


if __name__ == "__main__":
    main()
